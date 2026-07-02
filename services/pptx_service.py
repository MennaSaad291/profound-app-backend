"""
pptx_service.py — Matches Flutter preview exactly
==================================================
Layout per slide:
  - Colored header bar with white title text
  - Bullet points left column (60% or full width) — ▶ dot + bold headline + detail
  - Image right column (40%) when image_bytes_b64 present (fit=contain, no crop)
  - EXAMPLES box (green) with Real World / Practical / Industry / Analogy rows
  - CODE block (dark bg) with language label
  - DIAGRAM section: mermaid.ink rendered as image
  - QUIZ QUESTIONS section with numbered Qs and A/B/C/D options
  - DISCUSSION QUESTIONS bulleted list
  - ASSIGNMENTS bulleted list
  - Speaker notes → pptx notes pane
  - Slide number badge bottom-right
  - Title slide: full-bleed bg, decorative ovals, large title
"""
import io
import base64
import requests
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide dimensions (widescreen 13.333 x 7.5 in) ────────────────────────────
_W = Inches(13.333)
_H = Inches(7.5)

# ── Themes ────────────────────────────────────────────────────────────────────
THEMES = {
    "Modern Minimalist": {
        "bg":       (255,255,255), "header_bg":  (79,70,229),  "header_txt": (255,255,255),
        "accent":   (79,70,229),   "accent2":    (124,58,237),
        "title_txt":(30,27,75),    "hl_txt":     (30,27,75),   "det_txt":    (55,65,81),
        "dot":      (79,70,229),
        "ex_bg":    (236,253,245), "ex_lbl":     (6,95,70),    "ex_txt":     (6,95,70),
        "code_bg":  (30,30,46),    "code_txt":   (205,214,244),"code_lbl":   (139,213,202),
        "quiz_bg":  (245,243,255), "quiz_txt":   (60,7,100),   "quiz_opt":   (109,40,217),
        "disc_bg":  (245,243,255), "disc_txt":   (76,29,149),
        "asgn_bg":  (240,253,244), "asgn_txt":   (20,83,45),
        "diag_bg":  (238,242,255), "diag_lbl":   (79,70,229),
        "badge_bg": (79,70,229),   "badge_txt":  (255,255,255),
        "img_ph":   (224,231,255), "img_ph_txt": (79,70,229),
    },
    "Dark Mode Tech": {
        "bg":       (15,23,42),    "header_bg":  (30,40,64),   "header_txt": (56,189,248),
        "accent":   (56,189,248),  "accent2":    (6,182,212),
        "title_txt":(226,232,240), "hl_txt":     (226,232,240),"det_txt":    (148,163,184),
        "dot":      (56,189,248),
        "ex_bg":    (12,42,42),    "ex_lbl":     (52,211,153), "ex_txt":     (52,211,153),
        "code_bg":  (13,17,23),    "code_txt":   (137,220,235),"code_lbl":   (56,189,248),
        "quiz_bg":  (30,27,75),    "quiz_txt":   (196,181,253),"quiz_opt":   (167,139,250),
        "disc_bg":  (30,27,75),    "disc_txt":   (167,139,250),
        "asgn_bg":  (12,42,30),    "asgn_txt":   (52,211,153),
        "diag_bg":  (20,30,50),    "diag_lbl":   (56,189,248),
        "badge_bg": (56,189,248),  "badge_txt":  (15,23,42),
        "img_ph":   (30,64,79),    "img_ph_txt": (56,189,248),
    },
    "Classic Academic": {
        "bg":       (253,251,247), "header_bg":  (128,0,0),    "header_txt": (255,255,255),
        "accent":   (128,0,0),     "accent2":    (184,92,56),
        "title_txt":(59,10,10),    "hl_txt":     (59,10,10),   "det_txt":    (30,30,30),
        "dot":      (128,0,0),
        "ex_bg":    (240,247,238), "ex_lbl":     (26,71,42),   "ex_txt":     (26,71,42),
        "code_bg":  (30,30,46),    "code_txt":   (205,214,244),"code_lbl":   (184,92,56),
        "quiz_bg":  (253,245,221), "quiz_txt":   (112,58,0),   "quiz_opt":   (128,0,0),
        "disc_bg":  (253,245,221), "disc_txt":   (112,58,0),
        "asgn_bg":  (240,247,238), "asgn_txt":   (26,71,42),
        "diag_bg":  (248,244,238), "diag_lbl":   (128,0,0),
        "badge_bg": (128,0,0),     "badge_txt":  (255,255,255),
        "img_ph":   (245,230,208), "img_ph_txt": (128,0,0),
    },
    "Vibrant Creative": {
        "bg":       (255,255,255), "header_bg":  (249,115,22), "header_txt": (255,255,255),
        "accent":   (249,115,22),  "accent2":    (236,72,153),
        "title_txt":(67,20,7),     "hl_txt":     (67,20,7),    "det_txt":    (28,28,30),
        "dot":      (249,115,22),
        "ex_bg":    (253,242,248), "ex_lbl":     (112,26,117), "ex_txt":     (112,26,117),
        "code_bg":  (30,30,46),    "code_txt":   (205,214,244),"code_lbl":   (236,72,153),
        "quiz_bg":  (254,242,199), "quiz_txt":   (146,64,14),  "quiz_opt":   (180,50,10),
        "disc_bg":  (254,242,199), "disc_txt":   (146,64,14),
        "asgn_bg":  (253,242,248), "asgn_txt":   (112,26,117),
        "diag_bg":  (255,247,237), "diag_lbl":   (249,115,22),
        "badge_bg": (249,115,22),  "badge_txt":  (255,255,255),
        "img_ph":   (254,240,231), "img_ph_txt": (249,115,22),
    },
}


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def _rgb(t):
    return RGBColor(t[0], t[1], t[2])

def _rect(slide, l, t, w, h, color):
    """Add a filled rectangle with no border."""
    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(color)
    s.line.fill.background()
    return s

def _oval(slide, l, t, w, h, color):
    """Add a filled oval with no border."""
    s = slide.shapes.add_shape(9, int(l), int(t), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(color)
    s.line.fill.background()
    return s

def _tb(slide, l, t, w, h):
    """Add a textbox and return its text_frame."""
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    box.text_frame.word_wrap = True
    return box.text_frame

def _first_para(tf):
    """Return first paragraph if blank, else add new one."""
    if tf.paragraphs and tf.paragraphs[0].text == "":
        return tf.paragraphs[0]
    return tf.add_paragraph()

def _add_run(para, text, size, bold, color, italic=False):
    r = para.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = _rgb(color)
    return r

def _clean(s):
    import re
    return re.sub(r'\*+', '', str(s or '')).replace('__', '').strip()

def _ep(pt):
    """Extract (headline, detail) from a point dict or string."""
    if isinstance(pt, str):
        return _clean(pt), ""
    if isinstance(pt, dict):
        return (
            _clean(pt.get("headline", pt.get("title", ""))),
            _clean(pt.get("detail",   pt.get("explanation", ""))),
        )
    return _clean(str(pt)), ""

def _fit_image_bytes(img_bytes: bytes, box_w_px: int, box_h_px: int) -> bytes:
    """
    Scale image to fit within box_w x box_h (contain), no crop.
    Returns PNG bytes.
    """
    try:
        img = PILImage.open(io.BytesIO(img_bytes)).convert("RGBA")
        img.thumbnail((box_w_px, box_h_px), PILImage.LANCZOS)
        out = io.BytesIO()
        img.save(out, format="PNG")
        out.seek(0)
        return out.read()
    except Exception:
        return img_bytes

def _add_picture_fit(slide, img_bytes: bytes, l, t, w, h) -> bool:
    """Add picture scaled to fit (contain) within the given box."""
    try:
        # Scale to fit using PIL
        W_PX = int(w / Inches(1) * 96)
        H_PX = int(h / Inches(1) * 96)
        fitted = _fit_image_bytes(img_bytes, W_PX, H_PX)
        # Measure actual fitted size
        img = PILImage.open(io.BytesIO(fitted))
        iw, ih = img.size
        # Convert back to EMU
        scale = min(w / Inches(iw / 96), h / Inches(ih / 96))
        ew = int(Inches(iw / 96))
        eh = int(Inches(ih / 96))
        # Center in box
        el = int(l + (w - ew) / 2)
        et = int(t + (h - eh) / 2)
        slide.shapes.add_picture(io.BytesIO(fitted), el, et, ew, eh)
        return True
    except Exception as e:
        print(f"[PPTX] add_picture_fit failed: {e}")
        return False

# ── Mermaid diagram helper ────────────────────────────────────────────────────

def _fetch_mermaid_png(mermaid_code: str) -> bytes | None:
    """
    Fetch mermaid diagram as PNG from mermaid.ink with a strict 6s timeout.
    Returns None if fetch fails or times out — caller shows text fallback.
    """
    try:
        encoded = base64.urlsafe_b64encode(mermaid_code.encode("utf-8")).decode("ascii")
        url = f"https://mermaid.ink/img/{encoded}"
        resp = requests.get(url, timeout=6)
        if resp.status_code == 200 and len(resp.content) > 500:
            return resp.content
    except Exception as e:
        print(f"[PPTX] mermaid.ink: {e}")
    return None


# ── Title slide ───────────────────────────────────────────────────────────────

def _title_slide(slide, c, title: str, notes: str):
    # Full-bleed background
    _rect(slide, 0, 0, _W, _H, c["bg"])

    # Decorative ovals (top-right and bottom-right)
    _oval(slide, Inches(8.8),  Inches(-1.8), Inches(5.5), Inches(5.5), c["accent"])
    _oval(slide, Inches(10.2), Inches(4.0),  Inches(4.2), Inches(4.2), c["accent2"])

    # "UNIVERSITY LECTURE" tag
    tag = _tb(slide, Inches(1.1), Inches(1.2), Inches(7.5), Inches(0.45))
    p = _first_para(tag)
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, "UNIVERSITY LECTURE", 11, True, c["accent"])

    # Large title
    tf = _tb(slide, Inches(1.1), Inches(1.8), Inches(9.0), Inches(3.0))
    tf.word_wrap = True
    p = _first_para(tf)
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, title, 44, True, c["title_txt"])

    # Accent line under title
    _rect(slide, Inches(1.1), Inches(5.0), Inches(3.5), Inches(0.07), c["accent"])

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ── Section draw helpers ──────────────────────────────────────────────────────

def _draw_header(slide, c, title: str):
    """Draw colored header bar and return bottom Y."""
    _rect(slide, 0, 0, _W, Inches(0.9), c["header_bg"])
    ttf = _tb(slide, Inches(0.22), Inches(0.1), Inches(12.8), Inches(0.72))
    p = _first_para(ttf)
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, title, 22, True, c["header_txt"])
    return Inches(0.9)

def _draw_badge(slide, c, idx: int):
    """Draw slide number badge bottom-right."""
    bw, bh = Inches(0.55), Inches(0.30)
    bx = _W - bw - Inches(0.10)
    by = _H - bh - Inches(0.08)
    _rect(slide, bx, by, bw, bh, c["badge_bg"])
    btf = _tb(slide, bx + Inches(0.02), by + Inches(0.02), bw - Inches(0.04), bh - Inches(0.04))
    p = _first_para(btf)
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, str(idx), 9, True, c["badge_txt"])

def _draw_bullets(slide, c, pts: list, l, t, w, h):
    """Draw bullet points (▶ headline + detail) in a text box."""
    ptf = _tb(slide, l, t, w, h)
    for i, pt in enumerate(pts[:6]):
        hl, det = _ep(pt)
        if not hl:
            continue
        hp = _first_para(ptf) if i == 0 else ptf.add_paragraph()
        hp.space_before = Pt(10 if i > 0 else 2)
        hp.space_after  = Pt(1)
        # Bullet dot
        _add_run(hp, "▶  ", 10, True, c["dot"])
        # Headline
        _add_run(hp, hl, 17, True, c["hl_txt"])
        # Detail below
        if det:
            dp = ptf.add_paragraph()
            dp.space_before = Pt(1)
            dp.space_after  = Pt(1)
            _add_run(dp, "       " + det[:150], 12, False, c["det_txt"])

def _draw_examples(slide, c, example, practical, industry, analogy, l, t, w) -> float:
    """
    Draw green EXAMPLES box with sub-rows.
    Returns bottom Y after the box.
    """
    rows = [
        ("Real World:",  example),
        ("Practical:",   practical),
        ("Industry:",    industry),
        ("Analogy:",     analogy),
    ]
    rows = [(lbl, val) for lbl, val in rows if val and val.strip()]
    if not rows:
        return t

    # Estimate height: header row + each sub-row
    row_h = Inches(0.26)
    hdr_h = Inches(0.30)
    box_h = hdr_h + row_h * len(rows) + Inches(0.10)

    _rect(slide, l, t, w, box_h, c["ex_bg"])

    # "EXAMPLES" label header
    etf = _tb(slide, l + Inches(0.18), t + Inches(0.05), w - Inches(0.22), hdr_h)
    ep = _first_para(etf)
    ep.alignment = PP_ALIGN.LEFT
    _add_run(ep, "✦  EXAMPLES", 11, True, c["ex_lbl"])

    # Sub-rows
    for i, (lbl, val) in enumerate(rows):
        ry = t + hdr_h + i * row_h
        rtf = _tb(slide, l + Inches(0.28), ry, w - Inches(0.32), row_h)
        rp = _first_para(rtf)
        rp.alignment = PP_ALIGN.LEFT
        _add_run(rp, lbl + "  ", 11, True,  c["ex_lbl"])
        _add_run(rp, val[:120],  11, False, c["ex_txt"])

    return t + box_h + Inches(0.06)


def _draw_code(slide, c, code_ex: str, code_lang: str, l, t, w) -> float:
    """Draw dark code block. Returns bottom Y."""
    if not code_ex:
        return t

    lang_label = (code_lang or "CODE").upper()
    # Estimate height based on line count
    lines = code_ex.split("\n")
    line_count = min(len(lines), 12)
    box_h = Inches(0.38) + Inches(0.185) * line_count

    _rect(slide, l, t, w, box_h, c["code_bg"])

    ctf = _tb(slide, l + Inches(0.18), t + Inches(0.08), w - Inches(0.22), box_h - Inches(0.10))
    cp = _first_para(ctf)
    cp.alignment = PP_ALIGN.LEFT
    _add_run(cp, lang_label + "  ", 9, True, c["code_lbl"])
    _add_run(cp, "\n".join(lines[:12])[:320], 10, False, c["code_txt"])

    return t + box_h + Inches(0.06)


def _draw_diagram(slide, c, mermaid_code: str, l, t, w) -> float:
    """Embed mermaid diagram. Falls back to styled text box if network unavailable."""
    if not mermaid_code or not mermaid_code.strip():
        return t

    box_h = Inches(2.2)
    lbl_h = Inches(0.32)
    img_h = box_h - lbl_h - Inches(0.08)

    _rect(slide, l, t, w, box_h, c["diag_bg"])

    # Label bar
    dtf = _tb(slide, l + Inches(0.18), t + Inches(0.05), w - Inches(0.22), lbl_h)
    dp = _first_para(dtf)
    dp.alignment = PP_ALIGN.LEFT
    _add_run(dp, "⬡  DIAGRAM", 10, True, c["diag_lbl"])

    ix, iy, iw, ih = l + Inches(0.15), t + lbl_h, w - Inches(0.30), img_h

    # Try mermaid.ink with a thread timeout so we never block the export
    png = None
    try:
        import threading
        result = [None]
        def _fetch():
            result[0] = _fetch_mermaid_png(mermaid_code)
        th = threading.Thread(target=_fetch, daemon=True)
        th.start()
        th.join(timeout=8)   # max 8s for the entire diagram fetch
        png = result[0]
    except Exception:
        pass

    if png:
        if not _add_picture_fit(slide, png, ix, iy, iw, ih):
            png = None   # fallback to text

    if not png:
        # Render mermaid source as monospaced text in a dark box
        _rect(slide, ix, iy, iw, ih, c["code_bg"])
        ftf = _tb(slide, ix + Inches(0.10), iy + Inches(0.08), iw - Inches(0.15), ih - Inches(0.10))
        fp = _first_para(ftf)
        fp.alignment = PP_ALIGN.LEFT
        lines = mermaid_code.strip().split("\n")
        _add_run(fp, "\n".join(lines[:12]), 9, False, c["code_txt"])

    return t + box_h + Inches(0.06)

def _draw_quiz(slide, c, quiz_qs: list, l, t, w) -> float:
    """Draw QUIZ QUESTIONS section. Returns bottom Y."""
    if not quiz_qs:
        return t

    # Estimate height
    n_qs = min(len(quiz_qs), 4)
    q_h  = Inches(0.30)
    opt_h = Inches(0.24)
    hdr_h = Inches(0.32)
    box_h = hdr_h + n_qs * (q_h + 4 * opt_h) + Inches(0.12)
    # Cap to avoid overflow
    box_h = min(box_h, Inches(4.5))

    _rect(slide, l, t, w, box_h, c["quiz_bg"])

    # Header
    htf = _tb(slide, l + Inches(0.18), t + Inches(0.05), w - Inches(0.22), hdr_h)
    hp = _first_para(htf)
    hp.alignment = PP_ALIGN.LEFT
    _add_run(hp, "❓  QUIZ QUESTIONS", 11, True, c["quiz_txt"])

    cy = t + hdr_h
    for i, q in enumerate(quiz_qs[:n_qs]):
        if not isinstance(q, dict):
            continue
        question = _clean(q.get("question", ""))
        options  = q.get("options") or []

        # Question line
        qtf = _tb(slide, l + Inches(0.18), cy, w - Inches(0.22), q_h)
        qp = _first_para(qtf)
        qp.alignment = PP_ALIGN.LEFT
        _add_run(qp, f"{i+1}.  {question[:110]}", 12, True, c["quiz_txt"])
        cy += q_h

        # Options (A/B/C/D)
        for opt in options[:4]:
            opt_text = _clean(str(opt))
            otf = _tb(slide, l + Inches(0.50), cy, w - Inches(0.55), opt_h)
            op = _first_para(otf)
            op.alignment = PP_ALIGN.LEFT
            _add_run(op, opt_text[:100], 11, False, c["quiz_opt"])
            cy += opt_h

    return t + box_h + Inches(0.06)


def _draw_list_section(slide, c, items: list, label: str,
                       bg_key: str, txt_key: str, l, t, w) -> float:
    """Generic bulleted list section. Returns bottom Y."""
    if not items:
        return t

    n = min(len(items), 5)
    item_h = Inches(0.28)
    hdr_h  = Inches(0.32)
    box_h  = hdr_h + n * item_h + Inches(0.10)

    _rect(slide, l, t, w, box_h, c[bg_key])

    # Header
    htf = _tb(slide, l + Inches(0.18), t + Inches(0.05), w - Inches(0.22), hdr_h)
    hp = _first_para(htf)
    hp.alignment = PP_ALIGN.LEFT
    _add_run(hp, label, 11, True, c[txt_key])

    cy = t + hdr_h
    for item in items[:n]:
        text = _clean(str(item))
        itf  = _tb(slide, l + Inches(0.30), cy, w - Inches(0.35), item_h)
        ip   = _first_para(itf)
        ip.alignment = PP_ALIGN.LEFT
        _add_run(ip, "• " + text[:120], 11, False, c[txt_key])
        cy += item_h

    return t + box_h + Inches(0.06)

# ── Content slide ─────────────────────────────────────────────────────────────

def _content_slide(slide, c, idx: int, sd: dict):
    title     = _clean(sd.get("title",              "Slide"))
    pts       = sd.get("points",                    []) or []
    example   = _clean(sd.get("example",            ""))
    practical = _clean(sd.get("practical_example",  ""))
    industry  = _clean(sd.get("industry_example",   ""))
    analogy   = _clean(sd.get("analogy",            ""))
    code_ex   = _clean(sd.get("code_example",       ""))
    code_lang = _clean(sd.get("code_language",      "")).upper() or "CODE"
    mermaid   = (sd.get("diagram") or "").strip()
    notes     = _clean(sd.get("speaker_notes",      ""))
    quiz_qs   = sd.get("quiz_questions",            []) or []
    disc_qs   = sd.get("discussion_questions",      []) or []
    assigns   = sd.get("assignments",               []) or []

    # Resolve image bytes
    img_bytes = None
    b64 = sd.get("image_bytes_b64", "")
    if b64:
        try:
            img_bytes = base64.b64decode(b64)
        except Exception:
            pass

    has_img   = bool(img_bytes)
    has_ex    = bool(example or practical or industry or analogy)
    has_code  = bool(code_ex)
    has_diag  = bool(mermaid)
    has_quiz  = bool(quiz_qs)
    has_disc  = bool(disc_qs)
    has_asgn  = bool(assigns)

    # ── Background ────────────────────────────────────────────────────────────
    _rect(slide, 0, 0, _W, _H, c["bg"])

    # ── Header bar ────────────────────────────────────────────────────────────
    header_bottom = _draw_header(slide, c, title)

    # ── Slide number badge ────────────────────────────────────────────────────
    _draw_badge(slide, c, idx)

    # ── Layout: bullets left, image right ────────────────────────────────────
    MARGIN_L  = Inches(0.22)
    MARGIN_T  = header_bottom + Inches(0.10)
    AVAIL_W   = _W - Inches(0.44)   # total usable width

    if has_img:
        txt_w = Inches(7.4)
        img_w = Inches(5.3)
        img_l = MARGIN_L + txt_w + Inches(0.20)
    else:
        txt_w = AVAIL_W
        img_w = 0
        img_l = 0

    # Bullet area height: from MARGIN_T to near bottom (leave room for bottom sections)
    # Estimate how much bottom section needs
    bottom_h = Inches(0.10)
    if has_ex:    bottom_h += Inches(0.30) + Inches(0.26) * len([x for x in [example,practical,industry,analogy] if x and x.strip()]) + Inches(0.10) + Inches(0.06)
    if has_code:  bottom_h += Inches(0.38) + Inches(0.185) * min(len(code_ex.split("\n")),12) + Inches(0.06)
    if has_diag:  bottom_h += Inches(2.0) + Inches(0.06)
    if has_quiz:  bottom_h += min(Inches(0.32) + min(len(quiz_qs),4)*(Inches(0.30)+4*Inches(0.24))+Inches(0.12), Inches(4.5)) + Inches(0.06)
    if has_disc:  bottom_h += Inches(0.32) + min(len(disc_qs),5)*Inches(0.28) + Inches(0.10) + Inches(0.06)
    if has_asgn:  bottom_h += Inches(0.32) + min(len(assigns),5)*Inches(0.28) + Inches(0.10) + Inches(0.06)

    # Bullet area: from MARGIN_T to where bottom sections start
    bullet_avail = _H - MARGIN_T - bottom_h - Inches(0.40)
    bullet_avail = max(bullet_avail, Inches(1.5))  # always at least 1.5"

    # ── Bullet points (left column) ───────────────────────────────────────────
    _draw_bullets(slide, c, pts, MARGIN_L, MARGIN_T, txt_w, bullet_avail)

    # ── Image (right column) ─────────────────────────────────────────────────
    if has_img:
        img_h = min(bullet_avail, Inches(4.0))
        img_t = MARGIN_T + (bullet_avail - img_h) / 2
        if not _add_picture_fit(slide, img_bytes, img_l, img_t, img_w, img_h):
            _rect(slide, img_l, img_t, img_w, img_h, c["img_ph"])
            phtf = _tb(slide, img_l + Inches(0.2), img_t + img_h/2 - Inches(0.2), img_w - Inches(0.4), Inches(0.4))
            pp = _first_para(phtf)
            pp.alignment = PP_ALIGN.CENTER
            _add_run(pp, "Image", 12, False, c["img_ph_txt"])

    # ── Bottom sections (stacked vertically) ─────────────────────────────────
    cy = MARGIN_T + bullet_avail + Inches(0.10)
    full_w = _W - Inches(0.44)

    if has_code:
        cy = _draw_code(slide, c, code_ex, code_lang, MARGIN_L, cy, full_w)

    if has_ex:
        cy = _draw_examples(slide, c, example, practical, industry, analogy, MARGIN_L, cy, full_w)

    if has_diag:
        cy = _draw_diagram(slide, c, mermaid, MARGIN_L, cy, full_w)

    if has_quiz:
        cy = _draw_quiz(slide, c, quiz_qs, MARGIN_L, cy, full_w)

    if has_disc:
        cy = _draw_list_section(slide, c, disc_qs, "💬  DISCUSSION QUESTIONS",
                                "disc_bg", "disc_txt", MARGIN_L, cy, full_w)

    if has_asgn:
        cy = _draw_list_section(slide, c, assigns, "📝  ASSIGNMENTS",
                                "asgn_bg", "asgn_txt", MARGIN_L, cy, full_w)

    # ── Speaker notes ─────────────────────────────────────────────────────────
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ── Public entry point ────────────────────────────────────────────────────────

def create_pptx(data: dict) -> io.BytesIO:
    """
    Build and return a PPTX file matching the Flutter preview layout.

    data keys:
      theme      — one of the THEMES keys (default "Modern Minimalist")
      slides     — list of slide dicts from ai_lecture_generation_service
    """
    theme_name  = data.get("theme", "Modern Minimalist")
    c           = THEMES.get(theme_name, THEMES["Modern Minimalist"])
    slides_data = data.get("slides", [])

    if not slides_data:
        raise ValueError("No slides data provided")

    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    for idx, sd in enumerate(slides_data):
        if not isinstance(sd, dict):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        if idx == 0:
            _title_slide(
                slide, c,
                _clean(sd.get("title", "Lecture")),
                _clean(sd.get("speaker_notes", "")),
            )
        else:
            _content_slide(slide, c, idx + 1, sd)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
